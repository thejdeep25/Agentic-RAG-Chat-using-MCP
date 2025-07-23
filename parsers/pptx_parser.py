# This parser helps us extract all the text from PowerPoint slides.
from pptx import Presentation

def parse_pptx(file_path):
    """
    Opens a PPTX file and collects all the text from every slide and shape.
    Returns everything as one big string, separated by newlines.
    """
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
